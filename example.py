"""
Example usage of the PowerPoint Merger application.

This script demonstrates programmatic testing of the merge logic.
For normal usage, simply run: python merge_presentations.py
"""

from pptx import Presentation


def create_sample_presentations():
    """Create sample PowerPoint presentations for testing."""
    # Create first presentation
    prs1 = Presentation()
    slide = prs1.slides.add_slide(prs1.slide_layouts[0])
    title = slide.shapes.title
    title.text = "First Presentation"
    prs1.save('sample1.pptx')
    
    # Create second presentation
    prs2 = Presentation()
    slide = prs2.slides.add_slide(prs2.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Second Presentation"
    prs2.save('sample2.pptx')
    
    print("Sample presentations created: sample1.pptx and sample2.pptx")
    print("\nTo merge them:")
    print("1. Run: python merge_presentations.py")
    print("2. Enter: 2 (number of files)")
    print("3. Select or drag the sample files")
    print("4. Enter output name: merged_output")
    print("5. Reorder if desired and click 'Create New File'")


if __name__ == "__main__":
    create_sample_presentations()
