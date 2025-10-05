"""
PowerPoint Presentation Merger GUI Application

This script provides a graphical user interface for merging multiple PowerPoint
(.pptx) files into a single presentation file. The application guides users
through file selection, ordering, and naming, then launches the merged
presentation in slideshow mode.

Required libraries:
    - tkinter: GUI components
    - python-pptx: PowerPoint file handling
    - os, subprocess: File operations and launching PowerPoint
"""

import os
import subprocess
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD
from pptx import Presentation


class PowerPointMerger:
    """Main application class for merging PowerPoint presentations."""

    def __init__(self):
        """Initialize the application state."""
        self.num_files = 0
        self.selected_files = []
        self.output_filename = ""
        self.file_order = []

    def run(self):
        """Start the application with the first window."""
        self.show_number_of_files_window()

    def show_number_of_files_window(self):
        """Display Step 1: Window for entering number of files to merge."""
        window = tk.Tk()
        window.title("Step 1: Number of Files")
        window.geometry("400x150")

        # Label
        label = tk.Label(
            window,
            text="How many files should be merged?",
            font=("Arial", 12)
        )
        label.pack(pady=20)

        # Entry box
        entry = tk.Entry(window, font=("Arial", 12), width=10)
        entry.pack(pady=10)

        def on_next():
            """Handle Next button click."""
            try:
                # Validate input is a positive integer
                num = int(entry.get())
                if num <= 0:
                    raise ValueError("Number must be positive")
                self.num_files = num
                window.destroy()
                self.show_file_selection_window()
            except ValueError:
                messagebox.showerror(
                    "Invalid Input",
                    "Please enter a valid positive number."
                )

        # Next button
        next_btn = tk.Button(
            window,
            text="Next",
            command=on_next,
            font=("Arial", 12),
            width=10
        )
        next_btn.pack(pady=10)

        window.mainloop()

    def show_file_selection_window(self):
        """Display Step 2: Window for selecting files via drag-and-drop or dialog."""
        window = TkinterDnD.Tk()
        window.title("Step 2: Select Files")
        window.geometry("600x400")

        # Label
        label = tk.Label(
            window,
            text=f"Select {self.num_files} PowerPoint file(s):",
            font=("Arial", 12)
        )
        label.pack(pady=10)

        # Listbox for displaying files (with drag-and-drop support)
        listbox_frame = tk.Frame(window)
        listbox_frame.pack(pady=10, padx=20, fill=tk.BOTH, expand=True)

        scrollbar = tk.Scrollbar(listbox_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        listbox = tk.Listbox(
            listbox_frame,
            font=("Arial", 10),
            yscrollcommand=scrollbar.set
        )
        listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.config(command=listbox.yview)

        def on_drop(event):
            """Handle files dropped into the listbox."""
            # Parse dropped files (they come as a string)
            files = window.tk.splitlist(event.data)
            for file in files:
                # Only add .pptx files
                if file.endswith('.pptx') and file not in self.selected_files:
                    self.selected_files.append(file)
                    listbox.insert(tk.END, file)

        # Enable drag-and-drop on the listbox
        listbox.drop_target_register(DND_FILES)
        listbox.dnd_bind('<<Drop>>', on_drop)

        def add_files_from_disk():
            """Handle Add Files button click."""
            files = filedialog.askopenfilenames(
                title="Select PowerPoint Files",
                filetypes=[("PowerPoint Files", "*.pptx")]
            )
            for file in files:
                if file not in self.selected_files:
                    self.selected_files.append(file)
                    listbox.insert(tk.END, file)

        # Add Files button
        add_btn = tk.Button(
            window,
            text="Add Files from Disk",
            command=add_files_from_disk,
            font=("Arial", 12)
        )
        add_btn.pack(pady=5)

        def on_ok():
            """Handle OK button click."""
            if len(self.selected_files) != self.num_files:
                messagebox.showerror(
                    "Invalid Selection",
                    f"Please select exactly {self.num_files} file(s). "
                    f"You have selected {len(self.selected_files)} file(s)."
                )
            else:
                # Validate all files are valid .pptx files
                for file in self.selected_files:
                    if not os.path.exists(file):
                        messagebox.showerror(
                            "File Not Found",
                            f"File does not exist: {file}"
                        )
                        return
                    if not file.endswith('.pptx'):
                        messagebox.showerror(
                            "Invalid File",
                            f"File is not a .pptx file: {file}"
                        )
                        return
                window.destroy()
                self.show_filename_window()

        # OK button
        ok_btn = tk.Button(
            window,
            text="OK",
            command=on_ok,
            font=("Arial", 12),
            width=10
        )
        ok_btn.pack(pady=10)

        window.mainloop()

    def show_filename_window(self):
        """Display Step 3: Window for entering the output filename."""
        window = tk.Tk()
        window.title("New Filename")
        window.geometry("400x150")

        # Label
        label = tk.Label(
            window,
            text="New file name:",
            font=("Arial", 12)
        )
        label.pack(pady=20)

        # Entry box
        entry = tk.Entry(window, font=("Arial", 12), width=30)
        entry.pack(pady=10)

        def on_next():
            """Handle Next button click."""
            filename = entry.get().strip()
            if not filename:
                messagebox.showerror(
                    "Invalid Input",
                    "Please enter a filename."
                )
                return

            # Ensure .pptx extension
            if not filename.endswith('.pptx'):
                filename += '.pptx'

            self.output_filename = filename
            window.destroy()
            self.show_reorder_window()

        # Next button
        next_btn = tk.Button(
            window,
            text="Next",
            command=on_next,
            font=("Arial", 12),
            width=10
        )
        next_btn.pack(pady=10)

        window.mainloop()

    def show_reorder_window(self):
        """Display Step 4: Window for reordering files via drag-and-drop."""
        window = TkinterDnD.Tk()
        window.title("Step 3: Set Merge Order")
        window.geometry("600x400")

        # Label
        label = tk.Label(
            window,
            text="In what order should the presentations be merged? (Drag to reorder)",
            font=("Arial", 12),
            wraplength=550
        )
        label.pack(pady=10)

        # Initialize file order
        self.file_order = self.selected_files.copy()

        # Listbox for displaying and reordering files
        listbox_frame = tk.Frame(window)
        listbox_frame.pack(pady=10, padx=20, fill=tk.BOTH, expand=True)

        scrollbar = tk.Scrollbar(listbox_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        listbox = tk.Listbox(
            listbox_frame,
            font=("Arial", 10),
            yscrollcommand=scrollbar.set
        )
        listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.config(command=listbox.yview)

        # Populate listbox with filenames (not full paths for readability)
        for file in self.file_order:
            listbox.insert(tk.END, os.path.basename(file))

        # Variables for drag-and-drop reordering
        drag_data = {"index": None}

        def on_drag_start(event):
            """Handle drag start event."""
            widget = event.widget
            selection = widget.curselection()
            if selection:
                drag_data["index"] = selection[0]

        def on_drag_motion(event):
            """Handle drag motion event."""
            widget = event.widget
            index = widget.nearest(event.y)
            if index != drag_data["index"] and drag_data["index"] is not None:
                # Visual feedback: selection follows cursor
                widget.selection_clear(0, tk.END)
                widget.selection_set(index)

        def on_drop_reorder(event):
            """Handle drop event for reordering."""
            widget = event.widget
            start_index = drag_data["index"]
            if start_index is None:
                return

            drop_index = widget.nearest(event.y)

            if start_index != drop_index:
                # Reorder the file_order list
                item = self.file_order.pop(start_index)
                self.file_order.insert(drop_index, item)

                # Update the listbox
                listbox.delete(0, tk.END)
                for file in self.file_order:
                    listbox.insert(tk.END, os.path.basename(file))

                # Select the moved item
                listbox.selection_set(drop_index)

            drag_data["index"] = None

        # Bind drag-and-drop events
        listbox.bind('<Button-1>', on_drag_start)
        listbox.bind('<B1-Motion>', on_drag_motion)
        listbox.bind('<ButtonRelease-1>', on_drop_reorder)

        def on_create():
            """Handle Create New File button click."""
            window.destroy()
            self.merge_and_launch()

        # Create New File button
        create_btn = tk.Button(
            window,
            text="Create New File",
            command=on_create,
            font=("Arial", 12),
            width=15
        )
        create_btn.pack(pady=10)

        window.mainloop()

    def merge_and_launch(self):
        """Merge presentations and launch the slideshow."""
        try:
            # Create a new blank presentation
            merged_prs = Presentation()

            # Remove the default blank slide
            if len(merged_prs.slides) > 0:
                rId = merged_prs.slides._sldIdLst[0].rId
                merged_prs.part.drop_rel(rId)
                del merged_prs.slides._sldIdLst[0]

            # Iterate through source files in the specified order
            for file_path in self.file_order:
                try:
                    source_prs = Presentation(file_path)

                    # Copy each slide from the source presentation
                    for slide in source_prs.slides:
                        # Add a new slide with a blank layout
                        # Use the first available layout from the merged presentation
                        slide_layout = merged_prs.slide_layouts[0]
                        new_slide = merged_prs.slides.add_slide(slide_layout)

                        # Copy slide dimensions if needed
                        # Note: python-pptx handles this at the presentation level

                        # Copy all shapes from the source slide
                        for shape in slide.shapes:
                            self._copy_shape(shape, new_slide)

                except Exception as e:
                    messagebox.showerror(
                        "Error",
                        f"Failed to process file "
                        f"{os.path.basename(file_path)}: {str(e)}"
                    )
                    return

            # Save the merged presentation
            output_path = os.path.abspath(self.output_filename)
            merged_prs.save(output_path)

            # Launch PowerPoint in slideshow mode
            try:
                subprocess.Popen(
                    ['powerpnt.exe', '/s', output_path],
                    shell=True
                )
            except Exception as e:
                messagebox.showwarning(
                    "Launch Failed",
                    f"Presentation saved successfully to {output_path}, "
                    f"but failed to launch slideshow: {str(e)}"
                )

        except Exception as e:
            messagebox.showerror(
                "Error",
                f"Failed to merge presentations: {str(e)}"
            )

    def _copy_shape(self, source_shape, target_slide):
        """
        Copy a shape from source slide to target slide.

        Args:
            source_shape: The shape to copy from the source slide
            target_slide: The target slide to copy the shape to

        Note: This is a simplified implementation. Full shape copying
        with all properties is complex due to python-pptx limitations.
        """
        try:
            # Get the element of the source shape
            el = source_shape.element
            # Create a copy of the element
            new_el = el.__class__(el)
            # Add the copied element to the target slide
            target_slide.shapes._spTree.insert_element_before(
                new_el, 'p:extLst'
            )

        except Exception:
            # If direct copying fails, skip the shape
            # In production, you might want to handle specific shape types
            pass


def main():
    """Main entry point for the application."""
    app = PowerPointMerger()
    app.run()


if __name__ == "__main__":
    main()
